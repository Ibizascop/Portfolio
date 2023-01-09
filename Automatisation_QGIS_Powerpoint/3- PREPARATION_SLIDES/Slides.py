# -*- coding: utf-8 -*-


from pptx import Presentation
from pptx.util import Inches, Pt
from tqdm import tqdm
import support as sp
import pandas as pd
import os
import win32com.client
import time

def create_slide(group_name):
    """
    Fonction prenant en argument le nom d'un groupe hôtellier
    et créant un powerpoint ainsi qu'un PDF associé sur la base
    du modèle du ficchier Template_Global.pptx
    
    Etant donné que la fonction a besoin d'ouvrir et de fermer
    l'application POWERPOINT, il ne faut pas utiliser l'application
    pendant que cette fonction est en train de tourner
    
    Arguments:
        group_name: String: Le nom d'un groupe issu 
        de la colonne "Groups" du fichier Parc_Trip_2022.xlsx
        Par exemple : "ACCOR" 
        
    Renvoit: Powerpoint + PDF 
        Renvoit un powerpoint et un PDF avec le nom du groupe stockés
        respectivement dans les dossiers RESULTATS_SLIDES et PDF
    """
    # IMPORTS
    saving_name = group_name
    df_brand = sp.import_data(".\DATA\Parc_Trip_2022.xlsx")
    df_global = sp.import_data(".\DATA\Chiffres_Globaux_2022.xlsx")
    final_ppt = Presentation('.\PREPARATION_SLIDES\Template_Global.pptx')
    #For Group and each brand
    liste_brands = sp.list_brands(group_name,df_brand)
    limit = len(liste_brands) + 1
    for slide_number in range(0,limit)  : 
        #Group or Brand
        if slide_number != 0 : 
            global_group = False
            group_name = liste_brands[slide_number-1]
            segment = sp.get_segment(group_name,df_brand)
        else :
            global_group = True
            segment = None
            
        # Create variables
        if sp.domestic_hotel(group_name, df_brand,global_group) == False:
            country = sp.new_country_to_use(group_name,df_brand,global_group)
        else:
            country = sp.get_group_country(group_name, df_brand).title()
         
        global_ranking = sp.get_group_world_ranking(
             group_name, df_brand, global_group)            
        country_ranking = sp.get_group_country_ranking(
            group_name, country, df_brand, None, global_group)
        
        country_segment_ranking = sp.get_group_country_ranking(
            group_name, country, df_brand,segment,global_group)

        # Classement Monde global
        if global_ranking == 1 :
            template_world = 0
        else:
            template_world = 1
        # Classement pays global
        if country_ranking == 1:
            template = 2
        elif country_ranking == 2:
            template = 3
        elif country_ranking == 3:
            template = 4
        elif country_ranking >= 4 and not sp.is_country_last(group_name, country, df_brand,global_group,segment):
            template = 5
        else:
            template = 6
            
        # Classement pays par segments
        if country_segment_ranking == 1:
            segment_template = 7
        elif country_segment_ranking == 2:
            segment_template = 8
        elif country_segment_ranking == 3:
            segment_template = 9
        elif country_segment_ranking >= 4 and not sp.is_country_last(group_name, country, df_brand,global_group,segment):
            segment_template = 10
        else:
            segment_template = 11
        #Traitement slides   
        slide = final_ppt.slides[slide_number]
        shapes = slide.shapes

        # Change Images
        pictures = [shape for shape in shapes
                    if shape.shape_type == 13]
            
        # Change Map
        try:
            if slide_number == 0 :
                old_supply_map = pictures[2]
            else :
                old_supply_map = pictures[1]
            shapes.element.remove(old_supply_map.element)
            if len(liste_brands) == 1 :
                new_supply_map = '.\QGIS\CARTES\{}.png'.format(group_name.split("-")[0])
            else:
                new_supply_map = '.\QGIS\CARTES\{}.png'.format(group_name)
            slide.shapes.add_picture(new_supply_map, Inches(
                3.44/2.5), Inches(8.75/2.5), width=Inches(15.21/2.5))
        except FileNotFoundError:
            pass

        # Change Country icon
        if slide_number == 0 :
            old_country_icon = pictures[3]
        else:
            old_country_icon = pictures[2]
        shapes.element.remove(old_country_icon.element)
        if country =="Myanmar/Burma" :
            new_country_icon = '.\COUNTRY_ICONS\Myanmar.png'
        else:
            new_country_icon = '.\COUNTRY_ICONS\{}.png'.format(country)
        slide.shapes.add_picture(new_country_icon, Inches(
            19/2.5), Inches(2.38/2.5), width=Inches(2.64/2.5))

        # Change Global Piechart
        if slide_number == 0 :
            old_country_chart = pictures[4]
            shapes.element.remove(old_country_chart.element)
            sp.global_supply_graph(group_name,country, df_brand, df_global)
            new_country_chart = '.\GRAPHS\{}_global.png'.format(country)
            slide.shapes.add_picture(new_country_chart, Inches(
                18.5/2.5), Inches(11/2.5), width=Inches(7.54/2.5))

        # Change Branded Piechart
        if slide_number == 0 :
            old_group_chart = pictures[5]
            shapes.element.remove(old_group_chart.element)
            sp.branded_supply_graph(group_name, country,df_brand, df_global,global_group)
            new_group_chart = '.\GRAPHS\{}_branded.png'.format(group_name)
            slide.shapes.add_picture(new_group_chart, Inches(
                25.5/2.5), Inches(11/2.5), width=Inches(7.54/2.5))
            
        # Change Title Name
        titles = [shape for shape in shapes
                  if shape.shape_type == 14]

        title_textbox = titles[0]
        paragraph = title_textbox.text_frame.paragraphs[0]
        if slide_number != 0 : 
            line = "2022 RANKINGS : {} GLOBAL AND DOMESTIC* SUPPLY".format(
                group_name.split("-")[1].upper())
        else:
            line = "2022 RANKINGS : {} GLOBAL AND DOMESTIC* SUPPLY".format(
                group_name.upper())
        paragraph.runs[0].text = line

        # Modifier Blocs texte
        texts = [shape for shape in shapes
                 if shape.shape_type == 17]

        # Change Country Name
        country_textbox = texts[2]
        paragraph = country_textbox.text_frame.paragraphs[0]
        if country == "Bosnia and Herzegovina" :
            country_legend = "Bosnia"
        elif country =="United Republic of Tanzania" :
            country_legend = "Tanzania"
        elif country =="Myanmar/Burma" :
            country_legend = "Myanmar"
        else:
            country_legend =country
        if len(country_legend) >= 10:
            paragraph.font.size = Pt(15)
        if segment != None :
            country_segment_textbox = texts[4]
            segment_paragraph = country_segment_textbox.text_frame.paragraphs[0]
            if len(country_legend) >= 10:
                segment_paragraph.font.size = Pt(15)
            segment_table = {"SUPER-ECO":"BUDGET",
                             "ECO":"ECONOMY",
                             "MDG":"MIDSCALE",
                             "HDG":"UPSCALE",
                             "LUX":"LUXURY"}
            segment_text = segment_table[segment]
            paragraph.runs[0].text = "{} RANKING".format(country_legend.upper())
            segment_paragraph.runs[0].text = "{} RANKING".format(segment_text)
        else:
            paragraph.runs[0].text = "{} RANKING".format(country_legend.upper())

        # Modifier Tables
        tables = [shape for shape in shapes
                  if shape.shape_type == 19]

        # Update worldtable
        for i in [0,1] :
            if i != template_world :
                shapes.element.remove(tables[i].element)
        world_table = tables[template_world]
        world_table_data = sp.fill_global_table(group_name, df_brand,global_group)
        for i in range(1, 4):
            for j in range(4):
                cell = world_table.table.cell(i, j)
                paragraph = cell.text_frame.paragraphs[0]
                if len(paragraph.runs) > 1:
                    for bloc in paragraph.runs:
                        value = world_table_data["({},{})".format(i, j)]
                        if value is None :
                            value ="..."
                        try:
                            value = '{:,}'.format(
                                round(value))
                        except (ValueError, TypeError):
                            pass
                        bloc.text = value.upper()
                else:
                    run = paragraph.runs[0]
                    value = world_table_data["({},{})".format(i, j)]
                    if value is None :
                        value ="..."
                    try:
                        value = '{:,}'.format(round(value))
                    except (ValueError, TypeError):
                        pass
                    run.text = value.upper()

        # Update country table
        for i in range(2,7) :
            if i != template :
                shapes.element.remove(tables[i].element)
                country_table = tables[template]
    
        country_table_data = sp.fill_country_table(group_name, country, df_brand, template,global_group)
        # Change country name :
        cell = country_table.table.cell(0, 0)
        paragraph = cell.text_frame.paragraphs[0]
        run = paragraph.runs[0]
        run.text = "# {}".format(country)

        for i in range(1, 6):
            try:
                for j in range(4):
                    cell = country_table.table.cell(i, j)
                    paragraph = cell.text_frame.paragraphs[0]
                    if len(paragraph.runs) > 1:
                        for bloc in paragraph.runs:
                            value = country_table_data["({},{})".format(i, j)]
                            if value is None :
                                value ="..."
                            try:
                                value = '{:,}'.format(
                                    round(value))
                            except (ValueError, TypeError):
                                pass
                            bloc.text = value.upper()
                    else:
                        try:
                            run = paragraph.runs[0]
                            value = country_table_data["({},{})".format(i, j)]
                            if value is None :
                                value ="..."
                            value = '{:,}'.format(
                                round(value))
                        except (KeyError, IndexError, ValueError, TypeError):
                            pass
                        run.text = value.upper()
            except IndexError:
                break
            
        #Update country table by segments
        if slide_number != 0:
            for i in range(7,12) :
                if i != segment_template:
                    try :
                        shapes.element.remove(tables[i].element)
                    except Exception as e :
                        print("failed segments table == "+str(e))
            country_segment_table = tables[segment_template]
            
            country_segment_table_data = sp.fill_country_table_by_segments(group_name, country, df_brand, segment_template,segment,global_group)
            
            # Change country name :
            cell = country_segment_table.table.cell(0, 0)
            paragraph = cell.text_frame.paragraphs[0]
            run = paragraph.runs[0]
            run.text = "# {}".format(segment_text)
            
            for i in range(1, 6):
                try:
                    for j in range(4):
                        cell = country_segment_table.table.cell(i, j)
                        paragraph = cell.text_frame.paragraphs[0]
                        if len(paragraph.runs) > 1:
                            for bloc in paragraph.runs:
                                value = country_segment_table_data["({},{})".format(i, j)]
                                if value is None :
                                    value ="..."
                                try:
                                    value = '{:,}'.format(
                                        round(value))
                                except (ValueError, TypeError):
                                    pass
                                bloc.text = value.upper()
                        else:
                            try:
                                run = paragraph.runs[0]
                                value = country_segment_table_data["({},{})".format(i, j)]
                                if value is None :
                                    value ="..."
                                value = '{:,}'.format(
                                    round(value))
                            except (KeyError, IndexError, ValueError, TypeError):
                                pass
                            run.text = value.upper()
                except IndexError:
                    break
        else:
            pass
        
            
    #Remove unused slides
    for i in range(0, 41-len(liste_brands)):
        del final_ppt.slides._sldIdLst[len(liste_brands)+1]
    # Save as PPT
    final_ppt.save('.\RESULTAT_SLIDES\{}.pptx'.format(saving_name))

    #Save as PDF
    ppt_file_path = os.getcwd() + '.\RESULTAT_SLIDES\{}.pptx'.format(saving_name)
    powerpoint = win32com.client.Dispatch('Powerpoint.Application')
    deck = powerpoint.Presentations.Open(ppt_file_path)
    pdf_file_path = os.getcwd() + '.\PDF\{}.pdf'.format(saving_name)
    deck.SaveAs(pdf_file_path, 32)
    time.sleep(1)
    deck.Close()
    powerpoint.Quit()
    os.system('taskkill /F /IM POWERPNT.EXE')
 
def create_all_slides() :
    """
    Fonction permettant de lancer la création des slides et pdf
    pour les groupes n'étant pas listés dans le fichier done.txt,
    fichier indiquant quels groupes ont déjà été traités.
    
    La fonction peut être lancée directement en appelant le nom du fichier
    depuis un terminal Jupyter.
    
    Cette fonction fait appel à la fonction create_slide() et 
    la fait tourner sur l'ensemble des groupes.
    
    Etant donné que la fonction a besoin d'ouvrir et de fermer
    l'application POWERPOINT, il ne faut pas utiliser l'application
    pendant que cette fonction est en train de tourner
    
    Arguments:
        Aucun argument nécessaire
        
    Renvoit: Powerpoints + PDFs
        Renvoit les powerpoints et PDFs avec les noms des 
        différents groupes stockés respectivement dans les dossiers 
        RESULTATS_SLIDES et PDF. La fonction note également au fur etDATA à mesure
        quels groupes ont été traités.
    """
    #Lancer toutes les slides
    #Verifier les quelles ont déjà été faites
    try :
        with open(".\PREPARATION_SLIDES\done.txt","r") as flog :
            done_groups = flog.readlines()
            done_groups = [name.replace("\n","") for name in done_groups]
    except:
        done_groups = []
    groups = pd.read_excel(".\DATA\Parc_Trip_2022.xlsx")
    groups = groups.loc[groups['Rooms 2022'] != 0]
    groups = groups["Groups"].unique()
    groups = [group for group in groups if group not in done_groups]
    for group_name in tqdm(groups) :
        try :
            create_slide(group_name)
            with open(".\PREPARATION_SLIDES\done.txt","a") as file :
                print(group_name, file=file)
        except Exception as e:
            with open(".\PREPARATION_SLIDES\exception.txt","a") as file :
                print(group_name +" === "+str(e), file = file)
if __name__ == "__main__":
    create_all_slides()

        


